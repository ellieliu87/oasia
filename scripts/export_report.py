"""
export_report.py
================
Exports Dashboard, Portfolio Analytics, Security Analytics, and Attribution
to a .pptx file (PowerPoint format).

The .pptx file is directly compatible with:
  • Google Slides  — File → Open / Import slides (no API required)
  • Microsoft PowerPoint
  • LibreOffice Impress

No Google API credentials, no internet connection, and no authentication
are required. Everything runs locally.

What gets created (~15 slides)
-------------------------------
  Cover             — title, as-of date, key stats
  Dashboard         — KPI cards, sector donut, YTD bar, health radar
  Portfolio         — summary table, EVE stress chart, holdings table
  Security          — universe summary, pool detail + rate-shock table
  Attribution       — OAS / OAD / Yield / EVE waterfall charts + tables

Dependencies (beyond the project's existing packages)
------------------------------------------------------
    pip install python-pptx kaleido

Usage
-----
    python scripts/export_report.py
    python scripts/export_report.py --title "Q1 2026 Portfolio Review"
    python scripts/export_report.py --output reports/my_report.pptx
    python scripts/export_report.py --pool CC30_POOL_042
"""
from __future__ import annotations

import argparse
import io
import json
import sys
from datetime import date
from pathlib import Path

# ── Project root on sys.path ──────────────────────────────────────────────────
_ROOT = Path(__file__).parent.parent
sys.path.insert(0, str(_ROOT))


# ═════════════════════════════════════════════════════════════════════════════
# 1.  Colour palette  (matches the Oasia dark theme)
# ═════════════════════════════════════════════════════════════════════════════

from pptx.dml.color import RGBColor

BG       = RGBColor(0x0D, 0x11, 0x17)   # #0d1117  slide background
PANEL    = RGBColor(0x16, 0x1C, 0x26)   # #161c26  card / table row alt
ACCENT   = RGBColor(0x3A, 0x6F, 0xD1)   # #3a6fd1  headers, bars
TEXT_LT  = RGBColor(0xE6, 0xED, 0xF3)   # #e6edf3  primary text
TEXT_DIM = RGBColor(0x8B, 0x97, 0xA8)   # #8b97a8  secondary text
GREEN    = RGBColor(0x22, 0xC5, 0x5E)
RED      = RGBColor(0xEF, 0x44, 0x44)
AMBER    = RGBColor(0xDA, 0x9A, 0x20)


# ═════════════════════════════════════════════════════════════════════════════
# 2.  pptx helpers
# ═════════════════════════════════════════════════════════════════════════════

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN

# Slide size — matches Google Slides default (10" × 5.625", 16:9)
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)


def _new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _blank_slide(prs: Presentation):
    """Add a blank slide with the dark background colour."""
    layout = prs.slide_layouts[6]   # index 6 = blank
    slide  = prs.slides.add_slide(layout)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return slide


def _add_rect(slide, x_in: float, y_in: float, w_in: float, h_in: float,
              fill: RGBColor = ACCENT, line: RGBColor | None = None):
    from pptx.util import Inches
    shape = slide.shapes.add_shape(
        1,   # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()   # no border
    return shape


def _add_textbox(
    slide,
    text: str,
    x_in: float, y_in: float, w_in: float, h_in: float,
    font_size: float = 11,
    bold: bool  = False,
    italic: bool = False,
    color: RGBColor = TEXT_LT,
    align: str = "LEFT",
    word_wrap: bool = True,
) -> None:
    from pptx.util import Inches, Pt
    txb = slide.shapes.add_textbox(
        Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in)
    )
    tf = txb.text_frame
    tf.word_wrap = word_wrap
    p   = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(font_size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    p.alignment = {
        "LEFT":   PP_ALIGN.LEFT,
        "CENTER": PP_ALIGN.CENTER,
        "RIGHT":  PP_ALIGN.RIGHT,
    }.get(align, PP_ALIGN.LEFT)


def _add_image_bytes(slide, png_bytes: bytes,
                     x_in: float, y_in: float, w_in: float, h_in: float) -> None:
    slide.shapes.add_picture(
        io.BytesIO(png_bytes),
        Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in),
    )


def _add_table(
    slide,
    headers: list[str],
    rows: list[list[str]],
    x_in: float, y_in: float, w_in: float, h_in: float,
    font_size: float = 8,
) -> None:
    from pptx.util import Inches, Pt
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in),
    ).table

    # Even column widths
    col_w = Inches(w_in / n_cols)
    for col in tbl.columns:
        col.width = col_w

    def _style_cell(cell, text: str, bg: RGBColor,
                    txt_color: RGBColor = TEXT_LT,
                    bold: bool = False) -> None:
        cell.text = str(text)
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        tf = cell.text_frame
        tf.word_wrap = False
        p   = tf.paragraphs[0]
        run = p.runs[0] if p.runs else p.add_run()
        run.font.size  = Pt(font_size)
        run.font.bold  = bold
        run.font.color.rgb = txt_color

    # Header row
    for ci, hdr in enumerate(headers):
        _style_cell(tbl.cell(0, ci), hdr, ACCENT, TEXT_LT, bold=True)

    # Data rows
    for ri, row in enumerate(rows):
        bg = PANEL if ri % 2 == 0 else BG
        for ci, val in enumerate(row):
            _style_cell(tbl.cell(ri + 1, ci), str(val), bg)


# ── Slide layout helpers ──────────────────────────────────────────────────────

def _title_bar(slide, title: str) -> None:
    """Thin accent bar + title text at the top of a content slide."""
    _add_rect(slide, 0, 0, 10, 0.38, fill=PANEL)
    _add_textbox(slide, title, 0.15, 0.05, 9.7, 0.28,
                 font_size=13, bold=True, color=TEXT_LT)


def _section_slide(prs: Presentation, title: str) -> None:
    """Full-screen section divider with an accent side bar."""
    slide = _blank_slide(prs)
    _add_rect(slide, 0, 0, 0.12, 5.625, fill=ACCENT)
    _add_textbox(slide, title, 0.4, 2.2, 9.0, 0.9,
                 font_size=32, bold=True, color=TEXT_LT)


# ═════════════════════════════════════════════════════════════════════════════
# 3.  Data fetching  (uses the project's own tool layer)
# ═════════════════════════════════════════════════════════════════════════════

def _tool(name: str, args: dict = None) -> dict:
    from tool.registry import handle_tool_call
    raw = handle_tool_call(name, args or {})
    return json.loads(raw) if isinstance(raw, str) else raw


def _fetch_all_data(pool_id: str | None) -> dict:
    print("  portfolio summary …")
    summary    = _tool("get_portfolio_summary")
    print("  portfolio positions …")
    positions  = _tool("get_portfolio_positions")
    print("  EVE profile …")
    eve        = _tool("compute_eve_profile")
    print("  attribution …")
    attr_oas   = _tool("get_attribution", {"metric": "oas"})
    attr_oad   = _tool("get_attribution", {"metric": "oad"})
    attr_yield = _tool("get_attribution", {"metric": "yield"})
    attr_eve   = _tool("get_attribution", {"metric": "eve"})
    print("  sector allocation …")
    sectors    = _tool("get_sector_allocation")
    print("  top performers …")
    performers = _tool("get_top_performers")
    print("  portfolio health …")
    health     = _tool("get_portfolio_health")
    print("  universe summary …")
    universe   = _tool("get_universe_summary")

    sec_pool = pool_id or (
        (positions.get("positions") or [{}])[0].get("pool_id")
    )
    pool_details = shock_table = None
    if sec_pool:
        print(f"  pool details for {sec_pool} …")
        pool_details = _tool("get_pool_details",      {"pool_id": sec_pool})
        print(f"  rate shocks for {sec_pool} …")
        shock_table  = _tool("run_scenario_analysis", {"pool_ids": [sec_pool]})

    return dict(
        summary=summary, positions=positions, eve=eve,
        attr_oas=attr_oas, attr_oad=attr_oad, attr_yield=attr_yield, attr_eve=attr_eve,
        sectors=sectors, performers=performers, health=health,
        universe=universe, sec_pool=sec_pool,
        pool_details=pool_details, shock_table=shock_table,
    )


# ═════════════════════════════════════════════════════════════════════════════
# 4.  Plotly chart builders  →  PNG bytes (embedded directly, no upload)
# ═════════════════════════════════════════════════════════════════════════════

_W, _H = 1100, 560   # chart export resolution (px)

def _to_png(fig) -> bytes:
    import plotly.io as pio
    return pio.to_image(fig, format="png", width=_W, height=_H, scale=1.5)


def _chart_sector_donut(sectors_data: dict) -> bytes:
    import plotly.graph_objects as go
    items  = sectors_data.get("sectors", [])
    labels = [s["label"] for s in items]
    values = [s["mv"]    for s in items]
    colors = ["#3B6FD4", "#059669", "#D97706", "#a371f7",
              "#e11d48", "#0ea5e9", "#84cc16", "#f97316"]
    fig = go.Figure(go.Pie(
        labels=labels, values=values, hole=0.44,
        marker_colors=colors[:len(labels)],
        textinfo="label+percent",
    ))
    fig.update_layout(
        title="Sector Allocation by Market Value",
        paper_bgcolor="#0d1117", font_color="#e6edf3",
        margin=dict(t=55, b=15, l=15, r=15),
    )
    return _to_png(fig)


def _chart_ytd_bar(performers_data: dict) -> bytes:
    import plotly.graph_objects as go
    items  = performers_data.get("top", []) + performers_data.get("bottom", [])
    ids    = [p["pool_id"] for p in items]
    rets   = [p.get("ret_pct", 0) for p in items]
    colors = ["#22c55e" if r >= 0 else "#ef4444" for r in rets]
    fig = go.Figure(go.Bar(
        y=ids, x=rets, orientation="h",
        marker_color=colors,
        text=[f"{r:+.2f}%" for r in rets], textposition="outside",
    ))
    fig.update_layout(
        title="YTD Return — Top & Bottom Performers",
        xaxis_title="YTD Return (%)",
        paper_bgcolor="#0d1117", plot_bgcolor="#0d1117", font_color="#e6edf3",
        margin=dict(t=55, b=35, l=110, r=55),
    )
    fig.update_xaxes(gridcolor="#30363d", zerolinecolor="#58a6ff")
    return _to_png(fig)


def _chart_health_radar(health_data: dict) -> bytes:
    import plotly.graph_objects as go
    dims = health_data.get("sub_metrics", [])
    cats = [m["name"]  for m in dims]
    vals = [m["score"] for m in dims]
    cats_c = cats + [cats[0]]
    vals_c  = vals  + [vals[0]]
    fig = go.Figure()
    fig.add_trace(go.Scatterpolar(
        r=vals_c, theta=cats_c, fill="toself",
        fillcolor="rgba(58,111,209,0.28)", line_color="#3a6fd1", name="Portfolio",
    ))
    fig.add_trace(go.Scatterpolar(
        r=[5] * len(cats_c), theta=cats_c,
        line=dict(color="#58a6ff", dash="dash"), name="Reference",
    ))
    fig.update_layout(
        title=f"Portfolio Health  ·  Score: {health_data.get('health_score', 'N/A')}",
        polar=dict(
            radialaxis=dict(range=[0, 10], gridcolor="#30363d"),
            angularaxis=dict(gridcolor="#30363d"),
            bgcolor="#0d1117",
        ),
        paper_bgcolor="#0d1117", font_color="#e6edf3",
        margin=dict(t=55, b=25, l=35, r=35),
    )
    return _to_png(fig)


def _chart_eve(eve_data: dict) -> bytes:
    import plotly.graph_objects as go
    profile = eve_data.get("eve_profile", {})
    shocks  = sorted(profile.keys(),
                     key=lambda s: int(s.replace("bps","").replace("+","").replace("m","-")))
    pct  = [profile[s].get("pct_change", 0) for s in shocks]
    cols = ["#ef4444" if v < -5 else "#f97316" if v < 0 else "#22c55e" for v in pct]
    limit = eve_data.get("eve_limit_pct", -5)
    fig = go.Figure()
    fig.add_trace(go.Bar(x=shocks, y=pct, marker_color=cols,
                         text=[f"{v:.1f}%" for v in pct], textposition="outside"))
    fig.add_hline(y=limit, line_dash="dash", line_color="#ff6b6b",
                  annotation_text=f"Limit ({limit:.0f}%)", annotation_font_color="#ff6b6b")
    fig.update_layout(
        title="EVE Sensitivity — Rate Shock Scenarios",
        xaxis_title="Rate Shock", yaxis_title="EVE Change (%)",
        paper_bgcolor="#0d1117", plot_bgcolor="#0d1117", font_color="#e6edf3",
        margin=dict(t=55, b=35, l=55, r=35),
    )
    fig.update_xaxes(gridcolor="#30363d")
    fig.update_yaxes(gridcolor="#30363d", zerolinecolor="#58a6ff")
    return _to_png(fig)


def _chart_waterfall(attr_data: dict, title: str, unit: str, pos_color: str) -> bytes | None:
    import plotly.graph_objects as go
    drivers = attr_data.get("attribution", {})
    if not drivers:
        return None
    labels  = [d.replace("_", " ").title() for d in drivers] + ["Total"]
    values  = list(drivers.values())
    total   = sum(values)
    values.append(total)
    measure = ["relative"] * (len(values) - 1) + ["total"]
    fig = go.Figure(go.Waterfall(
        orientation="v", measure=measure,
        x=labels, y=values,
        connector={"line": {"color": "#30363d"}},
        increasing={"marker": {"color": pos_color}},
        decreasing={"marker": {"color": "#ef4444"}},
        totals={"marker":    {"color": "#58a6ff"}},
        text=[f"{v:+.2f}" for v in values], textposition="outside",
    ))
    fig.update_layout(
        title=f"{title}  ({unit})", yaxis_title=unit,
        paper_bgcolor="#0d1117", plot_bgcolor="#0d1117", font_color="#e6edf3",
        margin=dict(t=55, b=70, l=55, r=35),
    )
    fig.update_xaxes(gridcolor="#30363d")
    fig.update_yaxes(gridcolor="#30363d", zerolinecolor="#58a6ff")
    return _to_png(fig)


# ═════════════════════════════════════════════════════════════════════════════
# 5.  Slide builders  (one function per slide)
# ═════════════════════════════════════════════════════════════════════════════

def _safe(val, fmt=None, fallback="N/A") -> str:
    if val is None:
        return fallback
    try:
        return fmt.format(val) if fmt else str(val)
    except Exception:
        return str(val)


def slide_cover(prs: Presentation, title: str, as_of: str, summary: dict) -> None:
    slide = _blank_slide(prs)
    _add_rect(slide, 0, 0,    10, 0.1,  fill=ACCENT)
    _add_rect(slide, 0, 5.525, 10, 0.1, fill=ACCENT)
    _add_textbox(slide, title, 0.7, 1.5, 8.6, 1.0,
                 font_size=30, bold=True, color=TEXT_LT, align="CENTER")
    _add_textbox(slide, f"As of {as_of}  ·  Oasia MBS Analytics",
                 0.7, 2.55, 8.6, 0.4,
                 font_size=13, italic=True, color=TEXT_DIM, align="CENTER")
    stats = (
        f"MV: {_safe(summary.get('total_market_value'), '${:.1f}M')}   "
        f"OAS: {_safe(summary.get('weighted_oas_bps'), '{:.0f} bps')}   "
        f"OAD: {_safe(summary.get('weighted_oad_years'), '{:.2f} yr')}   "
        f"Yield: {_safe(summary.get('book_yield_pct'), '{:.2f}%')}"
    )
    _add_textbox(slide, stats, 0.7, 3.1, 8.6, 0.35,
                 font_size=11, color=TEXT_DIM, align="CENTER")


def slide_dashboard_kpis(prs: Presentation, data: dict) -> None:
    slide = _blank_slide(prs)
    _title_bar(slide, "Dashboard — Portfolio KPIs")
    s = data["summary"]
    breach = s.get("eve_breach", False)
    kpis = [
        ("Market Value",   _safe(s.get("total_market_value"),      "${:.1f}M")),
        ("Book Value",     _safe(s.get("total_book_value"),         "${:.1f}M")),
        ("Positions",      _safe(s.get("position_count"),           "{:,}")),
        ("OAS",            _safe(s.get("weighted_oas_bps"),         "{:.1f} bps")),
        ("OAD",            _safe(s.get("weighted_oad_years"),       "{:.2f} yr")),
        ("Book Yield",     _safe(s.get("book_yield_pct"),           "{:.2f}%")),
        ("Unrealized P&L", _safe(s.get("unrealized_pnl"),           "${:.0f}K")),
        ("EVE +200bp",     _safe(s.get("eve_up200_change_pct"),     "{:.1f}%")),
    ]
    card_w, card_h = 1.12, 0.9
    gap_x,  gap_y  = 0.15, 0.15
    start_x, start_y = 0.3, 0.55
    for i, (label, value) in enumerate(kpis):
        col = i % 4
        row = i // 4
        cx  = start_x + col * (card_w + gap_x)
        cy  = start_y + row * (card_h + gap_y)
        is_eve = "EVE" in label
        fill   = RGBColor(0x8C, 0x1A, 0x1A) if (is_eve and breach) else PANEL
        _add_rect(slide,  cx, cy, card_w, card_h, fill=fill)
        _add_textbox(slide, label, cx + 0.07, cy + 0.07, card_w - 0.14, 0.22,
                     font_size=7.5, color=TEXT_DIM)
        _add_textbox(slide, value, cx + 0.07, cy + 0.33, card_w - 0.14, 0.40,
                     font_size=18, bold=True, color=TEXT_LT)
    if breach:
        _add_textbox(slide, "⚠  EVE BREACH — rebalancing required",
                     0.3, 5.1, 6, 0.25, font_size=9, bold=True, color=RED)


def slide_sector_chart(prs: Presentation, data: dict) -> None:
    print("    rendering sector chart …")
    png   = _chart_sector_donut(data["sectors"])
    slide = _blank_slide(prs)
    _title_bar(slide, "Dashboard — Sector Allocation")
    _add_image_bytes(slide, png, 0.8, 0.45, 8.4, 4.9)


def slide_ytd_chart(prs: Presentation, data: dict) -> None:
    print("    rendering YTD chart …")
    png   = _chart_ytd_bar(data["performers"])
    slide = _blank_slide(prs)
    _title_bar(slide, "Dashboard — YTD Performance")
    _add_image_bytes(slide, png, 0.3, 0.45, 9.4, 4.9)


def slide_health_radar(prs: Presentation, data: dict) -> None:
    print("    rendering health radar …")
    png   = _chart_health_radar(data["health"])
    slide = _blank_slide(prs)
    _title_bar(slide, "Dashboard — Portfolio Health")
    _add_image_bytes(slide, png, 0.1, 0.42, 5.5, 4.7)
    dims = data["health"].get("sub_metrics", [])
    if dims:
        rows = [[m["name"], f"{m['score']:.1f}/10", m.get("desc", "")]
                for m in dims]
        _add_table(slide, ["Dimension", "Score", "Note"], rows,
                   5.7, 0.48, 4.1, min(len(rows) * 0.42 + 0.42, 4.8), font_size=8.5)


def slide_portfolio_summary(prs: Presentation, data: dict) -> None:
    s     = data["summary"]
    slide = _blank_slide(prs)
    _title_bar(slide, "Portfolio Analytics — Summary")
    rows = [
        ["Total Market Value",   _safe(s.get("total_market_value"),     "${:.2f}M")],
        ["Total Book Value",     _safe(s.get("total_book_value"),        "${:.2f}M")],
        ["Unrealized P&L",       _safe(s.get("unrealized_pnl"),          "${:.0f}K")],
        ["# Positions",          _safe(s.get("position_count"),          "{:,}")],
        ["Weighted OAS",         _safe(s.get("weighted_oas_bps"),        "{:.1f} bps")],
        ["Weighted OAD",         _safe(s.get("weighted_oad_years"),      "{:.2f} yr")],
        ["Weighted Convexity",   _safe(s.get("weighted_convexity"),      "{:.3f}")],
        ["Book Yield",           _safe(s.get("book_yield_pct"),          "{:.2f}%")],
        ["Annual Income",        _safe(s.get("annual_income"),           "${:.0f}K")],
        ["EVE Δ +200bp",         _safe(s.get("eve_up200_change_pct"),    "{:.1f}%")],
        ["EVE Limit",            _safe(s.get("eve_limit_pct"),           "{:.1f}%")],
    ]
    _add_table(slide, ["Metric", "Value"], rows, 0.2, 0.48, 4.5, 4.8, font_size=9)
    # Top-5 positions on the right
    pos_list = data["positions"].get("positions", [])[:6]
    if pos_list:
        pos_rows = [
            [p.get("pool_id", ""), p.get("product_type", ""),
             _safe(p.get("market_value"), "${:.1f}M"),
             _safe(p.get("oas_bps"),      "{:.0f}"),
             _safe(p.get("oad_years"),    "{:.2f}")]
            for p in pos_list
        ]
        _add_table(slide, ["Pool", "Type", "MV", "OAS", "OAD"],
                   pos_rows, 4.9, 0.48, 4.9, min(len(pos_rows) * 0.44 + 0.44, 4.8),
                   font_size=8.5)


def slide_eve_chart(prs: Presentation, data: dict) -> None:
    print("    rendering EVE chart …")
    png   = _chart_eve(data["eve"])
    slide = _blank_slide(prs)
    _title_bar(slide, "Portfolio Analytics — EVE Stress Test")
    _add_image_bytes(slide, png, 0.1, 0.42, 6.0, 4.8)
    profile = data["eve"].get("eve_profile", {})
    shocks  = sorted(profile.keys(),
                     key=lambda s: int(s.replace("bps","").replace("+","").replace("m","-")))
    rows = [
        [s,
         _safe(profile[s].get("pct_change"), "{:+.1f}%"),
         "⚠ BREACH" if profile[s].get("breach") else "OK"]
        for s in shocks
    ]
    _add_table(slide, ["Shock", "EVE Δ%", "Status"], rows,
               6.25, 0.48, 3.5, min(len(rows) * 0.44 + 0.44, 4.8), font_size=9)


def slide_holdings(prs: Presentation, data: dict) -> None:
    pos_list = data["positions"].get("positions", [])[:14]
    slide = _blank_slide(prs)
    _title_bar(slide, "Portfolio Analytics — Full Holdings")
    if not pos_list:
        _add_textbox(slide, "No positions available.", 0.3, 0.8, 9, 0.4)
        return
    rows = [
        [p.get("pool_id", ""),
         p.get("product_type", ""),
         _safe(p.get("coupon_pct"),          "{:.2f}%"),
         _safe(p.get("par_value"),           "${:.1f}M"),
         _safe(p.get("market_value"),        "${:.1f}M"),
         _safe(p.get("oas_bps"),             "{:.0f}"),
         _safe(p.get("oad_years"),           "{:.2f}"),
         _safe(p.get("unrealized_pnl_pct"),  "{:+.2f}%")]
        for p in pos_list
    ]
    _add_table(slide, ["Pool ID", "Type", "Coupon", "Par", "MV", "OAS", "OAD", "P&L%"],
               rows, 0.15, 0.48, 9.7, min(len(rows) * 0.34 + 0.38, 4.9), font_size=8)


def slide_universe_summary(prs: Presentation, data: dict) -> None:
    univ  = data["universe"]
    slide = _blank_slide(prs)
    _title_bar(slide, "Security Analytics — Universe Overview")
    by_prod = univ.get("by_product", [])
    if by_prod:
        rows = [
            [p.get("product_type", ""),
             str(p.get("count", "")),
             _safe(p.get("total_balance_bn"),  "{:.2f}B"),
             _safe(p.get("avg_coupon_pct"),    "{:.2f}%"),
             _safe(p.get("avg_oas_bps"),       "{:.0f}"),
             _safe(p.get("avg_oad_years"),     "{:.2f}"),
             _safe(p.get("avg_cpr_pct"),       "{:.1f}%"),
             _safe(p.get("avg_fico"),          "{:.0f}")]
            for p in by_prod
        ]
        _add_table(
            slide,
            ["Product", "Count", "Balance", "Avg Cpn", "Avg OAS", "Avg OAD", "Avg CPR", "Avg FICO"],
            rows, 0.15, 0.48, 9.7, min(len(rows) * 0.5 + 0.5, 4.8), font_size=9,
        )
    _add_textbox(
        slide,
        f"Total: {univ.get('total_pools','N/A')} pools  ·  "
        f"Balance: {_safe(univ.get('total_balance_bn'), '{:.2f}B')}",
        0.15, 5.2, 9.7, 0.25,
        font_size=8.5, italic=True, color=TEXT_DIM,
    )


def slide_pool_detail(prs: Presentation, data: dict) -> None:
    pool_id = data["sec_pool"]
    details = data["pool_details"]
    shocks  = data["shock_table"]
    if not pool_id or not details:
        return
    slide  = _blank_slide(prs)
    _title_bar(slide, f"Security Analytics — {pool_id}")
    static = details.get("static", {})
    info_rows = [
        ["Pool ID",      static.get("pool_id", pool_id)],
        ["Product",      static.get("product_type", "")],
        ["Coupon",       _safe(static.get("coupon_pct"),       "{:.2f}%")],
        ["WAC",          _safe(static.get("wac_pct"),          "{:.2f}%")],
        ["WALA",         _safe(static.get("wala_at_issue"),    "{:.0f} mo")],
        ["WAM",          _safe(static.get("original_wam"),     "{:.0f} mo")],
        ["FICO",         _safe(static.get("fico"),             "{:.0f}")],
        ["LTV",          _safe(static.get("ltv"),              "{:.1f}%")],
        ["Orig Balance", _safe(static.get("original_balance"), "${:.1f}M")],
    ]
    _add_table(slide, ["Characteristic", "Value"], info_rows,
               0.15, 0.48, 3.8, 4.5, font_size=9)
    if shocks and isinstance(shocks, dict):
        scenarios = shocks.get("scenarios", shocks.get("results", {}))
        # scenarios is a dict keyed by scenario name → unpack to rows
        if isinstance(scenarios, dict):
            sc_items = list(scenarios.items())[:10]
        else:
            sc_items = [(sc.get("scenario", ""), sc) for sc in scenarios[:10]]
        if sc_items:
            shock_rows = [
                [name,
                 _safe(sc.get("oas_bps"),    "{:.0f}"),
                 _safe(sc.get("oad_years"),  "{:.2f}"),
                 _safe(sc.get("convexity"),  "{:.3f}"),
                 _safe(sc.get("yield_pct"),  "{:.2f}%"),
                 _safe(sc.get("cpr_pct"),    "{:.1f}%")]
                for name, sc in sc_items
                if "error" not in sc
            ]
            _add_table(
                slide, ["Shock (bps)", "OAS", "OAD", "Convexity", "Yield", "CPR"],
                shock_rows, 4.15, 0.48, 5.7, min(len(shock_rows) * 0.42 + 0.42, 4.5),
                font_size=8.5,
            )


def slide_attribution(prs: Presentation, attr_data: dict,
                      title: str, unit: str, pos_color: str) -> None:
    print(f"    rendering {title} waterfall …")
    png = _chart_waterfall(attr_data, title, unit, pos_color)
    if png is None:
        return
    slide = _blank_slide(prs)
    _title_bar(slide, f"Attribution — {title}")
    _add_image_bytes(slide, png, 0.1, 0.42, 6.0, 4.8)
    drivers = attr_data.get("attribution", {})
    if drivers:
        total = sum(drivers.values())
        rows  = [[d.replace("_", " ").title(), f"{v:+.2f} {unit}"]
                 for d, v in drivers.items()]
        rows.append(["Total", f"{total:+.2f} {unit}"])
        _add_table(slide, ["Driver", "Value"], rows,
                   6.25, 0.48, 3.5, min(len(rows) * 0.44 + 0.44, 4.8), font_size=9)


# ═════════════════════════════════════════════════════════════════════════════
# 6.  Main orchestrator
# ═════════════════════════════════════════════════════════════════════════════

def build_report(title: str, pool_id: str | None, output_path: Path) -> Path:
    print("\n[1/3] Fetching data from Oasia …")
    data  = _fetch_all_data(pool_id)
    as_of = data["summary"].get("as_of_date", str(date.today()))

    print("\n[2/3] Building slides …")
    prs = _new_prs()

    # Cover
    print("  Cover")
    slide_cover(prs, title, as_of, data["summary"])

    # Dashboard
    print("  Section: Dashboard")
    _section_slide(prs, "Dashboard")
    slide_dashboard_kpis(prs, data)
    slide_sector_chart(prs, data)
    slide_ytd_chart(prs, data)
    slide_health_radar(prs, data)

    # Portfolio Analytics
    print("  Section: Portfolio Analytics")
    _section_slide(prs, "Portfolio Analytics")
    slide_portfolio_summary(prs, data)
    slide_eve_chart(prs, data)
    slide_holdings(prs, data)

    # Security Analytics
    print("  Section: Security Analytics")
    _section_slide(prs, "Security Analytics")
    slide_universe_summary(prs, data)
    slide_pool_detail(prs, data)

    # Attribution
    print("  Section: Attribution")
    _section_slide(prs, "Attribution Analysis")
    for attr_key, t, unit, color in [
        ("attr_oas",   "OAS Attribution",   "bps", "#16b3c5"),
        ("attr_oad",   "OAD Attribution",   "yrs", "#22c55e"),
        ("attr_yield", "Yield Attribution", "%",   "#da9a20"),
        ("attr_eve",   "EVE Attribution",   "$K",  "#16b3c5"),
    ]:
        slide_attribution(prs, data[attr_key], t, unit, color)

    print(f"\n[3/3] Saving to {output_path} …")
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


# ═════════════════════════════════════════════════════════════════════════════
# 7.  CLI entry point
# ═════════════════════════════════════════════════════════════════════════════

def main() -> None:
    parser = argparse.ArgumentParser(
        description="Export Oasia analytics to a .pptx file (Google Slides compatible)."
    )
    parser.add_argument(
        "--title",
        default=f"Oasia Portfolio Report — {date.today().strftime('%B %d, %Y')}",
        help="Presentation title",
    )
    parser.add_argument(
        "--output",
        default=None,
        help="Output file path (default: reports/oasia_report_<date>.pptx)",
    )
    parser.add_argument(
        "--pool",
        default=None,
        help="Pool ID to feature in the Security Analytics section",
    )
    args = parser.parse_args()

    output = Path(args.output) if args.output else (
        _ROOT / "reports" / f"oasia_report_{date.today().isoformat()}.pptx"
    )

    try:
        out = build_report(
            title=args.title,
            pool_id=args.pool,
            output_path=output,
        )
        print(f"\n✓  Report saved: {out}")
        print("   → Open in PowerPoint, or upload to Google Drive and open with Google Slides.")
    except Exception as exc:
        import traceback
        print(f"\nERROR: {exc}", file=sys.stderr)
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()
